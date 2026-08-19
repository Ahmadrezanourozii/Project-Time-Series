"""Build the Deliverable-3 presentation on the existing FAU deck template.

Starts from the Deliverable-2 file (so master, fonts and colours match),
removes its slides, and rebuilds the deck for a strict 15-minute talk.

Output: Mamba Results.AhmadrezaNourozi.pptx
"""

from __future__ import annotations

import copy
import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Emu, Inches, Pt

sys.path.insert(0, str(Path(__file__).resolve().parent.parent))

TEMPLATE = Path("Baseline Models.AhmadrezaNourozi.pptx")
OUT = Path("Mamba Results.AhmadrezaNourozi.pptx")
FINAL = Path("outputs/mamba/final")
FIGS = FINAL / "figures_report_png"

DARK_BLUE = RGBColor(0x00, 0x33, 0x66)
ORANGE = RGBColor(0xFF, 0x99, 0x33)
GREY = RGBColor(0x59, 0x59, 0x59)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

SLIDE_W, SLIDE_H = Emu(12192000), Emu(6858000)


def drop_all_slides(prs: Presentation) -> None:
    xml_slides = prs.slides._sldIdLst
    for sld in list(xml_slides):
        prs.part.drop_rel(sld.rId)
        xml_slides.remove(sld)


def layout(prs: Presentation, name: str):
    return next(l for l in prs.slide_layouts if l.name == name)


def add(prs: Presentation, layout_name: str, title: str | None = None):
    slide = prs.slides.add_slide(layout(prs, layout_name))
    for ph in list(slide.placeholders):
        idx = ph.placeholder_format.idx
        if idx == 0 and title is not None:
            ph.text_frame.text = title
            for p in ph.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size, r.font.bold, r.font.color.rgb = Pt(24), True, DARK_BLUE
        elif idx in (13, 14, 25, 26, 37, 39, 10, 11, 12):
            ph._element.getparent().remove(ph._element)  # unused placeholders
    return slide


def textbox(slide, x, y, w, h, lines, size=14, color=GREY, bold_first=False, bullet=True):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        text, lvl = (line if isinstance(line, tuple) else (line, 0))
        p.text = ("• " if bullet and lvl == 0 else ("– " if bullet else "")) + text
        p.level = lvl
        p.space_after = Pt(7)
        for r in p.runs:
            r.font.size = Pt(size - 1.5 * lvl)
            r.font.color.rgb = DARK_BLUE if (bold_first and i == 0) else color
            r.font.bold = bold_first and i == 0
    return box


def picture(slide, path: Path, x, y, w=None, h=None):
    kw = {}
    if w: kw["width"] = Inches(w)
    if h: kw["height"] = Inches(h)
    return slide.shapes.add_picture(str(path), Inches(x), Inches(y), **kw)


def table(slide, x, y, w, h, rows, col_widths=None, highlight_row=None, font_size=12):
    n_rows, n_cols = len(rows), len(rows[0])
    shape = slide.shapes.add_table(n_rows, n_cols, Inches(x), Inches(y), Inches(w), Inches(h))
    tbl = shape.table
    if col_widths:
        for i, cw in enumerate(col_widths):
            tbl.columns[i].width = Inches(cw)
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = tbl.cell(r, c)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.space_after = Pt(0)
                for run in p.runs:
                    run.font.size = Pt(font_size)
                    run.font.bold = (r == 0) or (highlight_row is not None and r == highlight_row)
                    if r == 0:
                        run.font.color.rgb = WHITE
                    elif highlight_row is not None and r == highlight_row:
                        run.font.color.rgb = DARK_BLUE
            if r == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = DARK_BLUE
            elif highlight_row is not None and r == highlight_row:
                cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(0xFF, 0xE8, 0xCC)
    return shape


def banner(slide, text, y=1.15, color=ORANGE):
    box = slide.shapes.add_textbox(Inches(0.55), Inches(y), Inches(12.2), Inches(0.42))
    tf = box.text_frame; tf.word_wrap = True
    tf.text = text
    for p in tf.paragraphs:
        for r in p.runs:
            r.font.size, r.font.bold, r.font.color.rgb = Pt(14), True, color
    return box


def main() -> None:
    mamba = json.loads((FINAL / "primary_amplitude30_3seed" / "results.json").read_text())
    base = json.loads((FINAL / "baseline_comparison.json").read_text())
    m = lambda foot, key: mamba["results"][foot]["subject_mean_prob"]["bootstrap"][key]
    b = lambda model, foot, key: base["subject_wise"][model][foot]["subject"]["bootstrap"][key]
    leak = lambda model, foot: base["leaky_window"][model][foot]["point"]["accuracy"] * 100

    prs = Presentation(str(TEMPLATE))
    drop_all_slides(prs)

    # 1 — title
    s = add(prs, "Titelbild")
    for ph in list(s.placeholders):
        ph._element.getparent().remove(ph._element)
    tb = s.shapes.add_textbox(Inches(0.9), Inches(2.3), Inches(11), Inches(2.2))
    tf = tb.text_frame; tf.word_wrap = True
    tf.text = "Deliverable 3: Classification with Mamba"
    tf.paragraphs[0].runs[0].font.size = Pt(40)
    tf.paragraphs[0].runs[0].font.bold = True
    tf.paragraphs[0].runs[0].font.color.rgb = DARK_BLUE
    for line, size, color in [
        ("Parkinson's disease detection from gait force signals", 20, GREY),
        ("with a state-space model — and why the protocol matters", 20, GREY),
        ("Ahmadreza Nourozi  ·  Project Time Series (SS 26)  ·  M.Sc. Artificial Intelligence", 14, GREY),
    ]:
        p = tf.add_paragraph(); p.text = line; p.space_before = Pt(10)
        p.runs[0].font.size = Pt(size); p.runs[0].font.color.rgb = color

    # 2 — where we are
    s = add(prs, "Textfolie", "Where we are")
    textbox(s, 0.6, 1.35, 6.0, 4.5, [
        "Deliverable 1 — data analysis, TSFresh features, PCA / t-SNE",
        "Deliverable 2 — baselines: Random Forest, SVM-RBF, 1D-ResNet, LSTM",
        "Deliverable 3 — classification with Mamba (state-space model)",
    ], size=15, bold_first=False)
    textbox(s, 6.9, 1.35, 5.8, 4.5, [
        "Two questions this talk answers",
        "Does a state-space model beat the baselines on this data?",
        "Why do published accuracies on this database differ so much?",
    ], size=15, bold_first=True)

    # 3 — data
    s = add(prs, "Textfolie", "Data and protocol")
    textbox(s, 0.6, 1.35, 6.1, 4.6, [
        "PhysioNet Gait in Parkinson's Disease (Ga, Ju, Si)",
        "100 subjects — 50 PD / 50 HC, balanced by design",
        "8 force sensors per foot + 2 totals, 100 Hz, ~2 min per subject",
        "First 5 s dropped (gait initiation); 5 s windows, 1 s step",
        "10 407 windows, 35–254 per subject",
    ], size=15)
    textbox(s, 6.9, 1.35, 5.8, 4.6, [
        "Five predefined folds — subject-wise",
        "60 train / 20 validation / 20 test subjects, class-balanced",
        "Test sets disjoint → pooling covers each subject exactly once",
        "z-score fitted on training windows only",
        "Hyper-parameters and threshold chosen on validation only",
    ], size=15, bold_first=True)

    # 4 — pipeline figure
    s = add(prs, "Nur Titel | weiß", "Pipeline")
    picture(s, FIGS / "fig1_pipeline.png", 0.5, 1.9, w=12.3)

    # 5 — model
    s = add(prs, "Textfolie", "Model: Mamba classifier")
    textbox(s, 0.6, 1.35, 6.1, 4.6, [
        "Mamba = selective state-space layer",
        "state transition depends on the input → keeps or forgets by content",
        "linear cost in sequence length (vs quadratic self-attention)",
        "Our classifier",
        "tokens → linear projection → 6 blocks (d = 256)",
        "mean ‖ max pooling over time → linear head → PD / HC",
    ], size=14, bold_first=True)
    textbox(s, 6.9, 1.35, 5.8, 4.6, [
        "Bidirectional variant — tested, not adopted",
        "a state-space scan is causal: early tokens see little context",
        "so we added a second mixer over the reversed sequence",
        "at matched size it gave no gain: 72% both, and the causal scan ranked better (AUC 0.87 vs 0.82)",
        "a negative result — mean-pooling may already carry that evidence",
    ], size=14, bold_first=True)

    # 6 — what to feed it
    s = add(prs, "Textfolie", "What we feed the model matters most")
    banner(s, "Raw signals made training collapse on 2 of 5 folds — feature-vector tokens fixed it")
    table(s, 0.75, 1.85, 7.4, 2.6, [
        ["Input representation", "Subject acc. (%)", "Worst fold val AUC"],
        ["Raw samples, 5 s windows", "69", "0.53"],
        ["Raw samples, 30 s / 60 s", "69 / 65", "0.63 / 0.59"],
        ["Stride-cycle tokens", "66", "0.71"],
        ["Amplitude ++ stride (fused)", "78", "0.71"],
        ["Amplitude statistics (used)", "80", "0.77"],
    ], col_widths=[3.4, 2.1, 1.9], highlight_row=5, font_size=12)
    textbox(s, 8.4, 1.9, 4.3, 4.2, [
        "One token = one second",
        "11 statistics × 18 channels = 198 dims",
        "mean, std, skew, kurtosis, RMS, ZCR, median, min, max, spectral energy & entropy",
        "The assignment allows raw signals and/or feature-vector sequences",
    ], size=13, bold_first=True)

    # 7 — main results table
    s = add(prs, "Nur Titel | weiß", "Results — subject-wise protocol (headline)")
    rows = [["Model", "Acc (%)", "Pre (%)", "Rec (%)", "F1", "Sen (%)", "Spe (%)", "AUC"]]
    for name, key in [("SVM-RBF", "svm_rbf"), ("1D-ResNet", "resnet"), ("LSTM", "lstm"),
                      ("Random Forest", "random_forest")]:
        rows.append([name, b(key, "both", "accuracy"), b(key, "both", "precision_w"),
                     b(key, "both", "recall_w"), b(key, "both", "f1_w"),
                     b(key, "both", "sensitivity"), b(key, "both", "specificity"),
                     b(key, "both", "auc")])
    rows.append(["Mamba (ours)", m("both", "accuracy"), m("both", "precision_w"),
                 m("both", "recall_w"), m("both", "f1_w"), m("both", "sensitivity"),
                 m("both", "specificity"), m("both", "auc")])
    table(s, 0.75, 1.75, 11.8, 3.0, rows, highlight_row=5, font_size=12)
    textbox(s, 0.75, 5.1, 11.8, 1.1, [
        "Both feet, 100 subjects pooled over the five disjoint test folds; mean ± std over 1000 bootstrap resamples.",
        "Precision / recall / F1 are class-weighted, so recall equals accuracy by construction. PD = positive class.",
    ], size=12, bullet=False)

    # 8 — per foot
    s = add(prs, "Nur Titel | weiß", "Results per foot configuration")
    picture(s, FIGS / "fig3_results.png", 0.6, 1.5, w=11.6)
    textbox(s, 0.75, 6.0, 11.8, 0.9, [
        "Mamba is best on the right foot (79%) and on both feet (78%); Random Forest keeps a small edge on the left foot (74% vs 70%).",
    ], size=12, bullet=False)

    # 9 — confusion matrices
    s = add(prs, "Nur Titel | weiß", "Confusion matrices — Mamba, subject level")
    picture(s, FIGS / "fig4_confusion.png", 1.4, 1.7, w=10.4)
    textbox(s, 0.75, 5.6, 11.8, 0.9, [
        "Both feet: 41 of 50 patients and 37 of 50 controls correct. Sensitivity (82%) exceeds specificity (74%) —",
        "the model rarely misses a patient, the preferable error profile for screening.",
    ], size=12, bullet=False)

    # 10 — the protocol issue
    s = add(prs, "Nur Titel | weiß", "Why published accuracies on this database disagree")
    picture(s, FIGS / "fig2_protocols.png", 0.9, 1.55, w=6.6)
    table(s, 8.0, 2.0, 4.6, 2.0, [
        ["Model", "Subject-wise", "Window-wise"],
        ["SVM-RBF", f"{b('svm_rbf','both','accuracy').split()[0]}%", f"{leak('svm_rbf','both'):.0f}%"],
        ["Random Forest", f"{b('random_forest','both','accuracy').split()[0]}%", f"{leak('random_forest','both'):.0f}%"],
    ], col_widths=[1.9, 1.4, 1.3], font_size=12)
    textbox(s, 8.0, 4.2, 4.6, 2.0, [
        "Same data, same features, same models, same metrics — only the split unit changes",
        "Windows of one recording are highly correlated: the model can recognise the person, not the disease",
        "We report subject-wise as our result; window-wise only as a warning",
    ], size=12.5, color=DARK_BLUE)

    # 11 — ablations
    s = add(prs, "Textfolie", "What actually moved the needle")
    table(s, 0.75, 1.6, 7.6, 3.1, [
        ["Change", "Subject acc. (%)"],
        ["Raw signal, base encoder", "69"],
        ["→ feature-vector tokens (same settings)", "76"],
        ["→ larger encoder (d = 128 → 256)", "80"],
        ["→ bidirectional scan (matched size)", "no gain"],
        ["→ 3-seed averaging (final, reported)", "78"],
    ], col_widths=[5.2, 2.4], highlight_row=5, font_size=12.5)
    textbox(s, 8.6, 1.65, 4.1, 4.3, [
        "Reading the table",
        "input representation and encoder capacity carried the gain",
        "bidirectionality: 72% either way at matched size, causal even ranked better (AUC 0.87 vs 0.82) — reported as a negative result",
        "seed averaging lowers the peak but is the honest number: single seeds ranged 72–80%",
        "models overfit by epoch 1–4 → warm-up + cosine, dropout ≤ 0.3",
    ], size=12, bold_first=True)

    # 12 — inference on unseen data
    s = add(prs, "Textfolie", "Ready for inference on a new cohort")
    textbox(s, 0.6, 1.35, 6.1, 4.6, [
        "What ships with the model",
        "one file per fold: weights + the per-channel mean/std of that fold's training subjects + its calibrated threshold",
        "new recording → trim 5 s → per-second tokens → 30 s windows → normalise with the stored statistics",
        "five fold models scored, averaged per window, then per recording",
    ], size=13.5, bold_first=True)
    textbox(s, 6.9, 1.35, 5.8, 4.6, [
        "Why the stored statistics matter",
        "recomputing mean/std on the new cohort would leak its own distribution into its predictions",
        "verified: windows rebuilt from a raw file at inference are numerically identical to the training windows",
        "one command:  python predict.py --bundle … --input <new .txt dir>",
    ], size=13.5, bold_first=True)

    # 13 — discussion
    s = add(prs, "Textfolie", "Discussion")
    textbox(s, 0.6, 1.35, 6.1, 4.6, [
        "Strengths",
        "subject-disjoint folds fixed in advance; test data used once",
        "normalisation, hyper-parameters and threshold from train/validation only",
        "identical evaluation code for every model, bootstrap dispersion reported",
    ], size=14, bold_first=True)
    textbox(s, 6.9, 1.35, 5.8, 4.6, [
        "Limitations",
        "100 subjects: one subject = one accuracy point, ±4 bootstrap spread",
        "only normal-walk recordings used (dual-task left out)",
        "stride detection by force threshold may misfire on severe gait impairment",
        "stride-cycle tokens alone underperform — their information is not yet fully exploited",
    ], size=14, bold_first=True)

    # 14 — conclusion
    s = add(prs, "Textfolie", "Conclusion")
    textbox(s, 0.6, 1.35, 12.1, 4.6, [
        "Mamba is the strongest model here under leakage-free evaluation: 78% accuracy, 0.78 weighted F1, 0.82 AUC on both feet (79% on the right foot)",
        "It beats the best classical baseline by 5 points and the best deep baseline by 12",
        "The gain came from the input representation (feature-vector tokens) and encoder capacity — not from bidirectionality, which showed no benefit at matched size",
        "The split unit changes the reported accuracy by ~20 points — an accuracy on this database is uninterpretable without it",
        "Next: more subjects (dual-task recordings, other cohorts), richer gait-cycle descriptors, calibrated per-subject uncertainty, UPDRS severity regression",
    ], size=15)

    # 15 — closing
    s = add(prs, "Schlussfolie")
    for ph in list(s.placeholders):
        ph._element.getparent().remove(ph._element)
    tb = s.shapes.add_textbox(Inches(0.9), Inches(2.6), Inches(11), Inches(1.6))
    tf = tb.text_frame
    tf.text = "Thank you!"
    tf.paragraphs[0].runs[0].font.size = Pt(40)
    tf.paragraphs[0].runs[0].font.bold = True
    tf.paragraphs[0].runs[0].font.color.rgb = DARK_BLUE
    p = tf.add_paragraph()
    p.text = "Code: github.com/Ahmadrezanourozii/Project-Time-Series   ·   Trained on Kaggle Tesla T4"
    p.space_before = Pt(14)
    p.runs[0].font.size = Pt(15); p.runs[0].font.color.rgb = GREY

    prs.save(str(OUT))
    print(f"wrote {OUT} with {len(prs.slides.__iter__.__self__._sldIdLst)} slides")


if __name__ == "__main__":
    main()
